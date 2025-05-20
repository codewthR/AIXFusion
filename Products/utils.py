# core/utils.py

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from nltk.tokenize import sent_tokenize

import nltk
nltk.download('punkt')


def analyze_text(text):
    sentences = sent_tokenize(text)
    slides = []
    chunk = []
    for i, sentence in enumerate(sentences, 1):
        chunk.append(sentence)
        if i % 3 == 0:
            slides.append(" ".join(chunk))
            chunk = []
    if chunk:
        slides.append(" ".join(chunk))
    return slides

def create_ppt_from_text(slides, theme="default"):
    prs = Presentation()
    if theme == "dark":
        background_color = (0, 0, 0)
        font_color = (255, 255, 255)
    else:
        background_color = None
        font_color = (0, 0, 0)

    for i, content in enumerate(slides):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title, body = slide.shapes.title, slide.placeholders[1]
        title.text = f"Slide {i+1}"
        body.text = content

        if theme == "dark":
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = prs.slide_master.background.fill.fore_color.rgb

    from io import BytesIO
    ppt_io = BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)
    return ppt_io
