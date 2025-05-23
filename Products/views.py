
from django.shortcuts import render
from docx import Document
import PyPDF2  # lighter than fitz
from pptx import Presentation
from pptx.util import Inches
import nltk


def presenting(response):
    return render(response,'product.html')



import os
import tempfile
from pptx import Presentation
from pptx.util import Inches
from django.http import HttpResponse
from django.views.decorators.csrf import csrf_exempt

import io
from pptx import Presentation
from django.http import HttpResponse
from django.views.decorators.csrf import csrf_exempt
from django.shortcuts import render
import google.generativeai as genai

# Configure Gemini API
genai.configure(api_key="AIzaSyBKsr-HEIlYsAxOJQXqiQ7goPY2WHPYLzk")  # Replace with your actual Gemini API key

@csrf_exempt
def generate_ppt(request):
    if request.method == 'POST':
        topic = request.POST.get('topic', 'Artificial Intelligence')

        # Step 1: Generate slide content using Gemini
        try:
            model = genai.GenerativeModel(model_name="gemini-2.0-flash")
            prompt = f"""Create a PowerPoint slide outline and content for a presentation on '{topic}'. 
Break the presentation into slides, and for each slide, provide:
- A title on the first line
- Followed by 3 to 5 bullet points (each starting with '-') on the next lines
Separate slides with a blank line."""

            response = model.generate_content(prompt)
            ai_text = response.text.strip()

            if not ai_text:
                return HttpResponse("Gemini did not return any content.")
        except Exception as e:
            return HttpResponse(f"Error generating content with Gemini: {str(e)}")

        # Step 2: Create PowerPoint
        prs = Presentation()
        slide_layout = prs.slide_layouts[1]  # Title and Content layout

        sections = ai_text.split('\n\n')
        for section in sections:
            lines = [line.strip() for line in section.strip().split('\n') if line.strip()]
            if not lines:
                continue
            title = lines[0]
            content = lines[1:] or ["(No content)"]

            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title
            slide.placeholders[1].text = "\n".join(content)

        # Step 3: Return PPTX file
        ppt_io = io.BytesIO()
        prs.save(ppt_io)
        ppt_io.seek(0)

        safe_filename = topic.replace(" ", "_").replace("/", "_")
        response = HttpResponse(
            ppt_io.read(),
            content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
        response['Content-Disposition'] = f'attachment; filename="{safe_filename}.pptx"'
        return response

    # GET request: show input form
    return render(request, 'writetext.html')













import nltk
nltk.download('punkt_tab')
# Ensure NLTK punkt tokenizer is available
nltk.download('punkt', quiet=True)
from nltk.tokenize import sent_tokenize


def home(request):
    return render(request, 'upload.html')


def extract_text_from_docx(path):
    doc = Document(path)
    return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])


def extract_text_from_pdf(path):
    text = ""
    with open(path, 'rb') as file:
        reader = PyPDF2.PdfReader(file)
        for page in reader.pages:
            text += page.extract_text() or ""
    return text


def split_text_into_chunks(text, max_chars=800):
    sentences = sent_tokenize(text)
    chunks, current = [], ""
    for sent in sentences:
        if len(current) + len(sent) <= max_chars:
            current += sent + " "
        else:
            chunks.append(current.strip())
            current = sent + " "
    if current:
        chunks.append(current.strip())
    return chunks


def basic_summarize(text):
    chunks = split_text_into_chunks(text)
    summary = []
    for chunk in chunks:
        lines = sent_tokenize(chunk)
        if len(lines) > 2:
            summary.append(lines[0] + " " + lines[1])
        else:
            summary.append(chunk)
    return summary


def convert_to_ppt(request):
    if request.method == 'POST' and request.FILES.get('input_file'):
        title = request.POST.get("Title", "Presentation Title")
        author = request.POST.get("author", "Author")
        input_file = request.FILES['input_file']
        ext = os.path.splitext(input_file.name)[1].lower()

        with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as temp:
            for chunk in input_file.chunks():
                temp.write(chunk)
            temp_path = temp.name

        # Extract text
        if ext in ['.doc', '.docx']:
            raw_text = extract_text_from_docx(temp_path)
        elif ext == '.pdf':
            raw_text = extract_text_from_pdf(temp_path)
        else:
            return HttpResponse("Unsupported file format.")

        # Summarize
        summaries = basic_summarize(raw_text)

        # PPT creation
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        content_slide_layout = prs.slide_layouts[1]

        # Title Slide
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = author

        # Content Slides
        for i, section in enumerate(summaries, 1):
            slide = prs.slides.add_slide(content_slide_layout)
            slide.shapes.title.text = f"Slide {i}"
            tf = slide.placeholders[1].text_frame
            for sent in sent_tokenize(section):
                if sent.strip():
                    tf.add_paragraph().text = "• " + sent.strip()

        # Serve PPT
        output_path = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx").name
        prs.save(output_path)
        with open(output_path, 'rb') as ppt_file:
            response = HttpResponse(
                ppt_file.read(),
                content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation'
            )
            response['Content-Disposition'] = 'attachment; filename="converted.pptx"'
            return response

    return HttpResponse("Invalid request.")
